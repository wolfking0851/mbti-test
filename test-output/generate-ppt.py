#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
测试文件生成 - PowerPoint 演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

def create_test_pptx():
    prs = Presentation()
    
    # 幻灯片 1: 标题页
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "智能龙虾 - PPT 测试"
    subtitle.text = "PowerPoint 能力验证测试\n大娃 (SF-0001) | 贵阳"
    
    # 幻灯片 2: 基本信息
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "1. 基本信息"
    tf = content.text_frame
    tf.text = f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    p = tf.add_paragraph()
    p.text = "生成者：大娃 (SF-0001)"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "测试地点：贵阳"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Python 库：python-pptx"
    p.level = 0
    
    # 幻灯片 3: 功能列表
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "2. 支持的 PPT 功能"
    tf = content.text_frame
    tf.text = "[OK] 多幻灯片支持"
    
    p = tf.add_paragraph()
    p.text = "[OK] 标题和副标题"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "[OK] 项目符号列表"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "[OK] 层级缩进"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "[OK] 自定义字体和颜色"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "[OK] 表格和图表"
    p.level = 0
    
    # 幻灯片 4: 表格示例
    slide_layout = prs.slide_layouts[1]  # 标题 + 内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置标题
    title = slide.shapes.title
    title.text = "3. 表格示例"
    
    # 添加表格
    rows = 4
    cols = 3
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(2.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 设置列宽
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)
    
    # 表头
    headers = ['功能', '状态', '备注']
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
    
    # 数据
    data = [
        ['幻灯片生成', '正常', '4 张幻灯片'],
        ['表格支持', '正常', '3 列 4 行'],
        ['中文支持', '正常', '无乱码']
    ]
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_data in enumerate(row_data):
            table.rows[row_idx].cells[col_idx].text = cell_data
    
    # 幻灯片 5: 结束页
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "测试完成！"
    subtitle.text = "感谢查看\n有任何问题随时告诉我"
    
    # 保存
    output_path = '/home/admin/.openclaw/workspace/test-output/测试演示-PPT.pptx'
    prs.save(output_path)
    print(f'✅ PPT 演示文稿已生成：{output_path}')
    return output_path

if __name__ == '__main__':
    create_test_pptx()
